from typing import List, Tuple
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
from dataclasses import dataclass

@dataclass
class Chunk:
    uid: str
    source_id: str  # page/slide number as string
    text: str

def extract_pdf_pages(path: str, start: int, end: int) -> List[Tuple[int,str]]:
    doc = fitz.open(path); n = doc.page_count
    start, end = max(1,start), min(end,n)
    out = []
    for i in range(start-1, end):
        page = doc.load_page(i)
        out.append((i+1, (page.get_text("text") or "").strip()))
    return out

def extract_pptx_slides(path: str, start: int, end: int) -> List[Tuple[int,str]]:
    prs = Presentation(path); n = len(prs.slides)
    start, end = max(1,start), min(end,n)
    out=[]
    for i in range(start-1, end):
        slide = prs.slides[i]
        texts=[]
        for sh in slide.shapes:
            if hasattr(sh, "text") and sh.text:
                texts.append(sh.text.strip())
        out.append((i+1, "\n".join(t for t in texts if t)))
    return out

def extract_docx_pages(path: str, start: int, end: int) -> List[Tuple[int,str]]:
    paras = [p.text.strip() for p in Document(path).paragraphs if p.text and p.text.strip()]
    if not paras: return []
    buckets, buf=[], ""
    for p in paras:
        if len(buf)+len(p)+1 > 3000:
            buckets.append(buf); buf = p
        else:
            buf = f"{buf}\n{p}" if buf else p
    if buf: buckets.append(buf)
    n = len(buckets); start, end = max(1,start), min(end,n)
    return [(i+1, buckets[i]) for i in range(start-1, end)]

def split_text(text: str, max_len=800, overlap=80) -> List[str]:
    if not text: return []
    words=text.split(); out=[]; cur=[]
    for w in words:
        cur.append(w)
        if sum(len(x)+1 for x in cur) > max_len:
            out.append(" ".join(cur[:-1]))
            cur = cur[-overlap:] if overlap < len(cur) else cur
    if cur: out.append(" ".join(cur))
    return out

def make_chunks(kind: str, items: List[Tuple[int,str]], corpus_id: str) -> List[Chunk]:
    chunks=[]
    for source_num, raw in items:
        for j, t in enumerate(split_text(raw, 800, 80)):
            chunks.append(Chunk(uid=f"{corpus_id}:{source_num}:{j}", source_id=str(source_num), text=t))
    return chunks
